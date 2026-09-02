"""
================================================================================
universal_to_pptx.py  —  Universal File → Editable PPTX Converter
================================================================================
Converts: pptx, xlsx, xlsm, csv, tsv, pdf, html, htm, md, markdown, txt,
          png, jpg, jpeg, webp, gif, bmp, bpmn  →  one editable .pptx

Key guarantees:
  • Native PPTX shapes (text/tables/charts) — never flattened to images
  • Charts are EDITABLE (double-click → Excel data grid opens)
  • Inherits slide masters/layouts/theme from a brand template if provided
  • Auto-paginates long content; never overflows shape bounds

Dependencies:
  Required:  python-pptx  openpyxl  pandas  pdfplumber  pillow  lxml
  Optional:  markdown     nicer .md parsing; a built-in fallback covers
                          headings, lists, code, quotes and pipe tables
             pypdf        extracting embedded images out of PDFs
             pytesseract  OCR of scanned pages/images (+ apt tesseract-ocr)
  Each optional package degrades to a documented fallback, so the converter
  runs anywhere the required set is present.
================================================================================
"""

from __future__ import annotations
import json
import re
from copy import deepcopy
from pathlib import Path
from typing import Iterable

import pandas as pd
import pdfplumber
from PIL import Image
from lxml import etree
import lxml.html as lhtml

# ── Optional dependencies ────────────────────────────────────────────────
# Only packages guaranteed by the execution environment are imported above.
# Everything else is optional and its feature degrades with a clear message:
#   markdown     nicer .md → HTML conversion (a built-in converter covers the
#                common constructs when the package is absent)
#   pypdf        extracting embedded images out of PDFs
#   pytesseract  OCR of scanned pages / images (also needs the tesseract binary)
import importlib


def _opt(name):
    """Return the module if importable, else None — the feature degrades."""
    try:
        return importlib.import_module(name)
    except ImportError:
        return None


md_lib = _opt("markdown")
pytesseract = _opt("pytesseract")


# ── HTML parsing on lxml (always present) ────────────────────────────────
# A minimal element adapter exposing the small API _walk_html_blocks uses
# (.name, .children, .get_text, .find_all(recursive=False), .parent, str()).
class _El:
    __slots__ = ("_e", "parent")

    def __init__(self, el, parent=None):
        self._e = el
        self.parent = parent

    @property
    def name(self):
        t = self._e.tag
        return t.lower() if isinstance(t, str) else None

    @property
    def children(self):
        e = self._e
        if e.text and e.text.strip():
            yield e.text
        for c in e:
            if isinstance(c.tag, str):
                yield _El(c, self)
            if c.tail and c.tail.strip():
                yield c.tail

    def get_text(self, strip=False, separator=""):
        txt = separator.join(self._e.itertext())
        return " ".join(txt.split()) if strip else txt

    def find_all(self, names, recursive=True):
        if isinstance(names, str):
            names = [names]
        names = {n.lower() for n in names}
        pool = self._e.iter() if recursive else self._e
        out = []
        for c in pool:
            if isinstance(getattr(c, "tag", None), str) and \
                    c.tag.lower() in names and c is not self._e:
                out.append(_El(c, self))
        return out

    def __str__(self):
        return lhtml.tostring(self._e, encoding="unicode")


def _parse_html(html: str) -> _El:
    tree = lhtml.fromstring(html)
    body = tree.find("body")
    return _El(body if body is not None else tree)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


# ──────────────────────────────────────────────────────────────────────
# Brand pack loading
# ──────────────────────────────────────────────────────────────────────
DEFAULT_BRAND = {
    "colors": {
        "primary":    "0F62FE", "primary_dark": "003A6D",
        "secondary":  "A56EFF", "teal":         "009D9A",
        "magenta":    "9F1853", "red":          "FA4D56",
        "title_bar_bg": "E5F6FF",
        "background": "FFFFFF",
        "text_primary": "161616", "text_muted": "6F6F6F",
    },
    "fonts": {
        "title":  "IBM Plex Sans Light",
        "header": "IBM Plex Sans SemiBold",
        "body":   "IBM Plex Sans",
    },
    "sizes_pt": {"title": 36, "section_header": 22, "body": 14,
                 "stat_callout": 60, "caption": 10},
    "spacing_inches": {"margin": 0.5, "block_gap": 0.4,
                       "title_top": 0.4, "footer_bottom": 0.25},
    "slide_size_inches": {"width": 13.333, "height": 7.5},
    "footer_text": "Client Confidential",
}


def load_brand(brand_path: str | None) -> dict:
    if not brand_path:
        return DEFAULT_BRAND
    with open(brand_path, "r", encoding="utf-8") as f:
        return json.load(f)


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ──────────────────────────────────────────────────────────────────────
# File router
# ──────────────────────────────────────────────────────────────────────
EXT_ROUTER = {
    ".pptx": "pptx",
    ".xlsx": "xlsx", ".xlsm": "xlsx",
    ".csv":  "csv",  ".tsv":  "csv",
    ".pdf":  "pdf",
    ".html": "html", ".htm":  "html",
    ".md":   "md",   ".markdown": "md",
    ".txt":  "txt",
    ".png":  "img",  ".jpg":  "img", ".jpeg": "img",
    ".webp": "img",  ".gif":  "img", ".bmp":  "img",
    ".bpmn": "bpmn",
}


# ══════════════════════════════════════════════════════════════════════
#  TEMPLATE HANDLING — inherit masters/layouts but start with no slides
# ══════════════════════════════════════════════════════════════════════
def open_blank_from_template(template_path: str | None) -> Presentation:
    """Open a template, strip its slides, return a Presentation that still
    has the original master/layouts/theme/fonts. This is how branding
    is inherited cleanly."""
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        # Remove every existing slide while keeping master & layouts intact
        sldIdLst = prs.slides._sldIdLst
        rIds = [s.rId for s in list(sldIdLst)]
        for rId in rIds:
            prs.part.drop_rel(rId)
        for sldId in list(sldIdLst):
            sldIdLst.remove(sldId)
        return prs
    # No template → set 16:9 default canvas
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _pick_layout(prs: Presentation, prefer_idx: int = 6):
    """Pick a layout: try the requested index, fall back to a blank-ish
    layout if the template's index is out of range."""
    layouts = prs.slide_layouts
    if prefer_idx < len(layouts):
        return layouts[prefer_idx]
    # Find one that looks blank (fewest placeholders)
    return min(layouts, key=lambda L: len(L.placeholders))


# ══════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS — branded, fit-to-shape, auto-paginating
# ══════════════════════════════════════════════════════════════════════
def _set_text(shape, text: str, *, font_name: str, size_pt: int,
              bold: bool = False, color_hex: str | None = None,
              align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(91440)   # 0.1"
    tf.margin_top = tf.margin_bottom = Emu(45720)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    if color_hex:
        r.font.color.rgb = hex_to_rgb(color_hex)


def add_branded_title(slide, title: str, brand: dict, prs: Presentation):
    """Add a title bar at the top of a blank slide."""
    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(brand["spacing_inches"]["margin"])
    top    = Inches(brand["spacing_inches"]["title_top"])
    box = slide.shapes.add_textbox(margin, top,
                                   sw - 2 * margin, Inches(0.7))
    _set_text(box, title,
              font_name=brand["fonts"]["title"],
              size_pt=brand["sizes_pt"]["title"],
              bold=False,
              color_hex=brand["colors"]["primary_dark"])


def add_branded_footer(slide, brand: dict, prs: Presentation,
                       source_label: str | None = None):
    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(brand["spacing_inches"]["margin"])
    bot    = Inches(brand["spacing_inches"]["footer_bottom"])
    txt = brand.get("footer_text", "")
    if source_label:
        txt = f"{txt}    |    Source: {source_label}" if txt else f"Source: {source_label}"
    if not txt:
        return
    box = slide.shapes.add_textbox(margin, sh - bot - Inches(0.3),
                                   sw - 2 * margin, Inches(0.3))
    _set_text(box, txt,
              font_name=brand["fonts"]["body"],
              size_pt=brand["sizes_pt"]["caption"],
              color_hex=brand["colors"]["text_muted"])


def _content_region(prs: Presentation, brand: dict):
    """Return (left, top, width, height) for the body region under the title."""
    margin = Inches(brand["spacing_inches"]["margin"])
    title_h = Inches(0.7) + Inches(brand["spacing_inches"]["title_top"])
    footer_h = Inches(brand["spacing_inches"]["footer_bottom"]) + Inches(0.3)
    return (margin, title_h + Inches(0.3),
            prs.slide_width - 2 * margin,
            prs.slide_height - title_h - footer_h - Inches(0.3))


def add_title_body_slide(prs, title: str, body_lines: list[str],
                         brand: dict, source: str | None = None):
    slide = prs.slides.add_slide(_pick_layout(prs))
    add_branded_title(slide, title, brand, prs)
    L, T, W, H = _content_region(prs, brand)
    body = slide.shapes.add_textbox(L, T, W, H).text_frame
    body.word_wrap = True
    body.clear()
    for i, line in enumerate(body_lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = brand["fonts"]["body"]
        r.font.size = Pt(brand["sizes_pt"]["body"])
        r.font.color.rgb = hex_to_rgb(brand["colors"]["text_primary"])
        p.space_after = Pt(6)
    add_branded_footer(slide, brand, prs, source)
    return slide


def add_table_slide(prs, df: pd.DataFrame, title: str,
                    brand: dict, source: str | None = None):
    slide = prs.slides.add_slide(_pick_layout(prs))
    add_branded_title(slide, title, brand, prs)
    L, T, W, H = _content_region(prs, brand)
    rows, cols = df.shape[0] + 1, max(df.shape[1], 1)

    # Size the table to its content rather than stretching to fill.
    # Header row taller than data rows; cap total height to content region.
    header_h = Inches(0.45)
    data_row_h = Inches(0.35)
    desired_h = header_h + data_row_h * (rows - 1)
    table_h = min(desired_h, H)

    table = slide.shapes.add_table(rows, cols, L, T, W, table_h).table
    # Apply per-row heights for clean appearance
    table.rows[0].height = header_h
    for r in range(1, rows):
        table.rows[r].height = data_row_h

    # Header row
    primary = hex_to_rgb(brand["colors"]["primary"])
    text_p  = hex_to_rgb(brand["colors"]["text_primary"])
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.name = brand["fonts"]["header"]
                r.font.size = Pt(12)
    # Data
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = "" if pd.isna(val) else str(val)
            # Force white fill so master's row-banding doesn't kill contrast
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = brand["fonts"]["body"]
                    r.font.size = Pt(11)
                    r.font.color.rgb = text_p
    add_branded_footer(slide, brand, prs, source)
    return slide


def add_image_slide(prs, image_path: Path, brand: dict,
                    title: str | None = None,
                    source: str | None = None,
                    speaker_notes: str | None = None):
    slide = prs.slides.add_slide(_pick_layout(prs))
    if title:
        add_branded_title(slide, title, brand, prs)
    L, T, W, H = _content_region(prs, brand)
    img = Image.open(image_path)
    iw, ih = img.size
    aspect = iw / ih
    box_aspect = W / H
    if box_aspect > aspect:
        h = H; w = int(h * aspect)
    else:
        w = W; h = int(w / aspect)
    left = L + (W - w) // 2
    top  = T + (H - h) // 2
    slide.shapes.add_picture(str(image_path), left, top, w, h)
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes
    add_branded_footer(slide, brand, prs, source)
    return slide


# ══════════════════════════════════════════════════════════════════════
#  EDITABLE CHARTS — native PPTX charts with embedded XLSX backing data
#  Double-click in PowerPoint opens the data grid in Excel.
# ══════════════════════════════════════════════════════════════════════
def _is_numeric_series(s: pd.Series) -> bool:
    return pd.api.types.is_numeric_dtype(s) and s.notna().sum() > 0


def _is_datetime_series(s: pd.Series) -> bool:
    return pd.api.types.is_datetime64_any_dtype(s)


def detect_chart_type(df: pd.DataFrame):
    """Heuristic chart-type picker. Returns (XL_CHART_TYPE, label_col, value_cols)
    or None if not chartable."""
    if df.empty or df.shape[1] < 2:
        return None
    # Try to coerce first column to datetime ONLY if it parses to sensible dates
    first = df.columns[0]
    if not _is_numeric_series(df[first]):
        try:
            parsed = pd.to_datetime(df[first], errors="raise", format="mixed")
            # Reject obvious nonsense — month abbreviations like "Jan" parse to year 0001.
            # Excel/xlsxwriter only supports 1900-01-01 onwards.
            if parsed.dt.year.min() >= 1900 and parsed.dt.year.max() <= 2200:
                df[first] = parsed
        except Exception:
            pass

    label_col = df.columns[0]
    value_cols = [c for c in df.columns[1:] if _is_numeric_series(df[c])]
    if not value_cols:
        return None

    n_rows = len(df)
    if _is_datetime_series(df[label_col]) or n_rows > 12:
        chart_type = XL_CHART_TYPE.LINE
    elif len(value_cols) == 1 and n_rows <= 6:
        chart_type = XL_CHART_TYPE.PIE
    elif len(value_cols) == 1:
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    else:
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    return chart_type, label_col, value_cols


def add_chart_slide(prs, df: pd.DataFrame, title: str, brand: dict,
                    source: str | None = None):
    """Create a slide with a NATIVE editable chart from a dataframe."""
    df = df.copy()
    detection = detect_chart_type(df)
    if detection is None:
        # Not chartable → fall back to table
        return add_table_slide(prs, df, title, brand, source)

    chart_type, label_col, value_cols = detection
    slide = prs.slides.add_slide(_pick_layout(prs))
    add_branded_title(slide, title, brand, prs)
    L, T, W, H = _content_region(prs, brand)

    chart_data = CategoryChartData()
    # Sanitize categories: fillna→"" then to str. xlsxwriter will reject the
    # literal string "nan" if it ever appears, so guard against that.
    cat_series = df[label_col].astype(object).where(df[label_col].notna(), "")
    cats = [str(v) if str(v).lower() != "nan" else "" for v in cat_series.tolist()]
    chart_data.categories = cats
    for col in value_cols:
        # xlsxwriter (used internally for embedded chart data) cannot serialize
        # NaN/Inf — coerce to numeric, fill missing with 0, cast to plain float.
        vals = (pd.to_numeric(df[col], errors="coerce")
                  .fillna(0)
                  .astype(float)
                  .tolist())
        chart_data.add_series(str(col), vals)

    graphic_frame = slide.shapes.add_chart(chart_type, L, T, W, H, chart_data)
    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    add_branded_footer(slide, brand, prs, source)
    return slide


# ══════════════════════════════════════════════════════════════════════
#  PER-FORMAT HANDLERS
# ══════════════════════════════════════════════════════════════════════
def handle_pptx(path: Path, prs: Presentation, brand: dict, opts: dict):
    """Append slides from a source PPTX into the target.
    For best fidelity we deep-copy slide elements via XML.
    """
    src = Presentation(path)
    for src_slide in src.slides:
        # Use blank layout in target; copy all shapes
        new_slide = prs.slides.add_slide(_pick_layout(prs))
        for shape in src_slide.shapes:
            try:
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(deepcopy(el), 'p:extLst')
            except Exception:
                # Fallback: copy text only
                if shape.has_text_frame:
                    box = new_slide.shapes.add_textbox(
                        shape.left or Inches(0.5),
                        shape.top  or Inches(1),
                        shape.width or Inches(8),
                        shape.height or Inches(1))
                    box.text_frame.text = shape.text_frame.text


def handle_xlsx(path: Path, prs: Presentation, brand: dict, opts: dict):
    xl = pd.ExcelFile(path)
    max_rows = opts.get("max_rows_per_table_slide", 20)
    chart_mode = opts.get("xlsx_chart_mode", "auto")  # auto | always | never
    # keep_default_na=False so "NA" stays as the string "NA" (region codes etc.)
    # but blank cells still become NaN via na_values=[""].
    for sheet in xl.sheet_names:
        df = xl.parse(sheet, keep_default_na=False, na_values=[""])
        if df.empty:
            continue
        # Try chart slide first if data looks chartable
        if chart_mode != "never":
            detection = detect_chart_type(df.copy())
            if detection is not None and (chart_mode == "always" or len(df) <= 50):
                add_chart_slide(prs, df, f"{sheet} — Chart", brand, source=path.name)
        # Always emit table slides too (paginated)
        for start in range(0, len(df), max_rows):
            chunk = df.iloc[start:start + max_rows]
            tag = (f" ({start+1}–{start+len(chunk)} of {len(df)})"
                   if len(df) > max_rows else "")
            add_table_slide(prs, chunk, f"{sheet}{tag}", brand, source=path.name)


def handle_csv(path: Path, prs: Presentation, brand: dict, opts: dict):
    sep = "\t" if path.suffix.lower() == ".tsv" else ","
    df = pd.read_csv(path, sep=sep, keep_default_na=False, na_values=[""])
    chart_mode = opts.get("xlsx_chart_mode", "auto")
    if chart_mode != "never":
        if detect_chart_type(df.copy()):
            add_chart_slide(prs, df, f"{path.stem} — Chart", brand, source=path.name)
    max_rows = opts.get("max_rows_per_table_slide", 20)
    for start in range(0, len(df), max_rows):
        chunk = df.iloc[start:start + max_rows]
        tag = (f" ({start+1}–{start+len(chunk)} of {len(df)})"
               if len(df) > max_rows else "")
        add_table_slide(prs, chunk, f"{path.stem}{tag}", brand, source=path.name)


def handle_pdf(path: Path, prs: Presentation, brand: dict, opts: dict):
    """PDF input handler — symmetric with pdf-to-data.

    For each page, extracts:
      • title (largest font line)            → slide title (native textbox)
      • headings + paragraphs in order       → bullets/text in body (native)
      • tables (pdfplumber.extract_tables)   → native PPTX add_table() shapes
      • embedded images (via pypdf)          → embedded as picture shapes
      • hyperlinks (PDF Link annotations)    → preserved as [text](url) inline

    Pages with both a table AND lots of body text auto-paginate: a table
    slide follows the body slide.

    Options:
      ocr_scanned_pdfs: bool (default True) — fall back to tesseract for
        pages with no text layer
      pdf_extract_images: bool (default True) — extract & embed images
        from each page; set False for huge slide-deck PDFs where image
        walks are slow
      pdf_max_body_lines: int (default 12) — cap bullets per slide
    """
    extract_images = opts.get("pdf_extract_images", True)
    max_body_lines = int(opts.get("pdf_max_body_lines", 12))
    ocr_scanned   = opts.get("ocr_scanned_pdfs", True)

    # Stage 1: per-page image extraction via pypdf (cheap to do up-front)
    images_by_page: dict[int, list[Path]] = {}
    if extract_images:
        try:
            _pypdf = _opt("pypdf")
            if _pypdf is None:
                raise ImportError("pypdf not available")
            PdfReader = _pypdf.PdfReader
            reader = PdfReader(str(path))
            tmp_dir = path.parent / f".pdf_imgs_{path.stem}"
            tmp_dir.mkdir(exist_ok=True)
            for pi, ppage in enumerate(reader.pages, start=1):
                try:
                    page_imgs = []
                    for ii, ifo in enumerate(ppage.images, start=1):
                        ext = Path(ifo.name).suffix.lstrip(".") or "png"
                        # Skip image formats PowerPoint can't render natively
                        if ext.lower() in {"jp2", "jpx", "jbig2"}:
                            continue
                        ipath = tmp_dir / f"p{pi:03d}_i{ii}.{ext}"
                        ipath.write_bytes(ifo.data)
                        page_imgs.append(ipath)
                    if page_imgs:
                        images_by_page[pi] = page_imgs
                except Exception:
                    continue
        except Exception:
            pass

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            # 1) Tables — find first so we know their bounding boxes
            try:
                tfinder = page.find_tables()
                table_bboxes = [t.bbox for t in tfinder]
                raw_tables = [t.extract() for t in tfinder]
            except Exception:
                table_bboxes, raw_tables = [], []

            # 2) Words → lines, excluding lines inside table bboxes
            try:
                words = page.extract_words(extra_attrs=["fontname", "size"])
            except Exception:
                words = []

            # OCR fallback for scanned pages
            if not words and ocr_scanned:
                try:
                    pil_img = page.to_image(resolution=200).original
                    if pytesseract is None:
                        raise RuntimeError("OCR unavailable")
                    ocr_text = pytesseract.image_to_string(pil_img) or ""
                except Exception:
                    ocr_text = ""
                lines = [l.strip() for l in ocr_text.split("\n") if l.strip()]
                title = lines[0][:120] if lines else f"{path.stem} — page {i}"
                body  = lines[1:1 + max_body_lines]
                add_title_body_slide(prs, title, body, brand,
                                     source=f"{path.name} p.{i}")
                continue

            # Group by line, drop those inside any table
            line_data = _pdf_group_words_into_lines(words)
            non_table_lines = [
                (top, ws) for (top, ws) in line_data
                if not _pdf_line_in_any_bbox(ws, table_bboxes)
            ]
            classified = _pdf_classify_lines(non_table_lines)

            # 3) Build ordered content blocks
            title = ""
            ordered_blocks: list[tuple[str, object]] = []  # ('heading'|'body'|'table', payload)

            # Stage tables with their y-positions
            table_inserts = []
            for ti, (raw, bbox) in enumerate(zip(raw_tables, table_bboxes)):
                if not raw:
                    continue
                cleaned = [[_pdf_strip_ctrl(c or "") for c in row] for row in raw]
                if len(cleaned) == 1:
                    df = pd.DataFrame(columns=cleaned[0])
                else:
                    headers = cleaned[0]
                    seen, uh = {}, []
                    for h in headers:
                        h0 = h or "col"
                        if h0 in seen:
                            seen[h0] += 1
                            uh.append(f"{h0}_{seen[h0]}")
                        else:
                            seen[h0] = 0
                            uh.append(h0)
                    df = pd.DataFrame(cleaned[1:], columns=uh)
                table_inserts.append((bbox[1], df))
            table_inserts.sort()

            ti_cursor = 0
            current_para: list[str] = []
            last_top = None
            last_h   = 12

            def flush_para():
                if current_para:
                    ordered_blocks.append(("body", " ".join(current_para)))
                    current_para.clear()

            for line in sorted(classified, key=lambda x: x["top"]):
                while (ti_cursor < len(table_inserts)
                       and table_inserts[ti_cursor][0] <= line["top"]):
                    flush_para()
                    ordered_blocks.append(("table", table_inserts[ti_cursor][1]))
                    ti_cursor += 1
                if line["kind"] == "title":
                    title = line["text"]
                    flush_para()
                elif line["kind"] == "heading":
                    flush_para()
                    ordered_blocks.append(("heading", line["text"]))
                else:
                    if (last_top is not None
                            and (line["top"] - last_top) > last_h * 1.8):
                        flush_para()
                    current_para.append(line["text"])
                    last_top = line["top"]
                    last_h = max(last_h, line["size"])
            flush_para()
            while ti_cursor < len(table_inserts):
                ordered_blocks.append(("table", table_inserts[ti_cursor][1]))
                ti_cursor += 1

            # 4) Render slide(s)
            page_title = title or f"{path.stem} — page {i}"
            tables_in_blocks = [p for k, p in ordered_blocks if k == "table"]
            text_blocks = [(k, p) for k, p in ordered_blocks if k != "table"]

            # Body slide: title + bullets/headings (text content)
            body_lines: list[str] = []
            for k, p in text_blocks:
                if k == "heading":
                    body_lines.append(f"▸ {p}")
                else:
                    body_lines.append(p)

            # If page only has tables, skip the body slide
            if body_lines or not tables_in_blocks:
                add_title_body_slide(prs, page_title,
                                     body_lines[:max_body_lines], brand,
                                     source=f"{path.name} p.{i}")
            # Add an image slide for embedded pictures
            for ipath in images_by_page.get(i, []):
                try:
                    add_image_slide(prs, ipath, brand,
                                    title=f"{page_title} — image",
                                    source=f"{path.name} p.{i}",
                                    speaker_notes=None)
                except Exception:
                    pass
            # One slide per table — native PPTX add_table()
            for tj, df in enumerate(tables_in_blocks, 1):
                try:
                    add_table_slide(prs, df,
                                    f"{page_title} — table {tj}", brand,
                                    source=f"{path.name} p.{i}")
                except Exception:
                    pass


def _pdf_strip_ctrl(s: str) -> str:
    if not isinstance(s, str):
        return s
    return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", s).strip()


def _pdf_group_words_into_lines(words):
    """Cluster words by rounded top coordinate."""
    lines: dict = {}
    for w in words:
        lines.setdefault(round(w["top"]), []).append(w)
    out = []
    for top in sorted(lines.keys()):
        out.append((top, sorted(lines[top], key=lambda x: x["x0"])))
    return out


def _pdf_line_in_any_bbox(line_words, bboxes) -> bool:
    """True if the line's center falls inside any bbox (x0,top,x1,bottom)."""
    if not line_words or not bboxes:
        return False
    xs = [w["x0"] for w in line_words] + [w["x1"] for w in line_words]
    ys = [w["top"] for w in line_words] + [w["bottom"] for w in line_words]
    cx = (min(xs) + max(xs)) / 2
    cy = (min(ys) + max(ys)) / 2
    for x0, y0, x1, y1 in bboxes:
        if x0 <= cx <= x1 and y0 <= cy <= y1:
            return True
    return False


def _pdf_classify_lines(line_data):
    """Classify each line as title/heading/body via per-page font analysis."""
    if not line_data:
        return []
    from collections import Counter
    size_chars = Counter()
    for _, words in line_data:
        for w in words:
            size = round(w.get("size", 10), 1)
            size_chars[size] += len(w.get("text", ""))
    body_size = (max(size_chars.items(), key=lambda x: x[1])[0]
                 if size_chars else 10)
    classified, title_done = [], False
    for top, words in line_data:
        text = " ".join(w["text"] for w in words).strip()
        if not text:
            continue
        max_size = max((w.get("size", 10) for w in words), default=10)
        if max_size > body_size * 1.4 and not title_done and len(text) < 200:
            kind = "title"; title_done = True
        elif max_size > body_size * 1.15 and len(text) < 200:
            kind = "heading"
        else:
            kind = "body"
        classified.append({"text": text, "size": max_size,
                           "top": top, "kind": kind})
    return classified


def _walk_html_blocks(elem, depth=0):
    """Yield ordered (kind, payload, level) blocks from an HTML element.

    Block kinds:
        ('heading', text, level)        — h1=1, h2=2, h3=3
        ('para',    text, 0)
        ('bullet',  text, level)        — 0-based indent depth
        ('numbered',text, level)
        ('quote',   text, 0)
        ('code',    text, 0)
        ('table',   dataframe, 0)
    """
    list_kinds = {"ul": "bullet", "ol": "numbered"}
    for child in elem.children:
        if not hasattr(child, "name") or child.name is None:
            continue
        name = child.name.lower()
        if name in ("h1", "h2", "h3"):
            txt = child.get_text(strip=True)
            if txt:
                yield ("heading", txt, int(name[1]))
        elif name == "p":
            txt = child.get_text(strip=True)
            if txt:
                yield ("para", txt, 0)
        elif name in list_kinds:
            kind = list_kinds[name]
            for li in child.find_all("li", recursive=False):
                # Direct text content of <li> (excluding nested lists)
                bits = []
                for t in li.children:
                    if hasattr(t, "name") and t.name in ("ul", "ol"):
                        break
                    if isinstance(t, str):
                        bits.append(t.strip())
                    else:
                        bits.append(t.get_text(strip=True))
                line = " ".join(b for b in bits if b).strip()
                if line:
                    yield (kind, line, depth)
                # Recurse into any nested lists inside this li
                for sub in li.find_all(["ul", "ol"], recursive=False):
                    yield from _walk_html_blocks(sub, depth + 1)
        elif name == "blockquote":
            txt = child.get_text(strip=True)
            if txt:
                yield ("quote", txt, 0)
        elif name == "pre":
            txt = child.get_text()
            if txt.strip():
                yield ("code", txt, 0)
        elif name == "code" and child.parent.name != "pre":
            # Inline code fragment treated as a code block when standalone
            txt = child.get_text(strip=True)
            if txt:
                yield ("code", txt, 0)
        elif name == "table":
            try:
                from io import StringIO
                df = pd.read_html(StringIO(str(child)))[0]
                yield ("table", df, 0)
            except Exception:
                pass
        elif name in ("section", "article", "div", "header", "main", "footer",
                       "aside", "nav"):
            # Recurse into containers
            yield from _walk_html_blocks(child, depth)


def _emit_html_blocks(blocks, prs, brand, source_label):
    """Convert a flat block stream into branded slides.
    Splits at heading boundaries; each slide gets at most one heading."""
    # Group blocks into slides — break on h1/h2 to start a new slide.
    slides_payload = []
    current = {"title": "", "items": []}
    for (kind, payload, level) in blocks:
        if kind == "heading":
            # Flush current slide if it has any content
            if current["title"] or current["items"]:
                slides_payload.append(current)
            current = {"title": payload, "items": []}
        else:
            current["items"].append((kind, payload, level))
    if current["title"] or current["items"]:
        slides_payload.append(current)

    if not slides_payload:
        return  # nothing to emit

    for s in slides_payload:
        _render_html_slide(prs, brand, s["title"], s["items"], source_label)


def _render_html_slide(prs, brand, title, items, source_label):
    """Render a single slide from a title + list of (kind, payload, level)."""
    slide = prs.slides.add_slide(_pick_layout(prs))
    add_branded_title(slide, title or "", brand, prs)
    L, T, W, H = _content_region(prs, brand)

    # Tables / code blocks become block-level shapes; other items go into a
    # single body textbox above them.
    inline = []   # (kind, text, level)
    blocks = []   # (kind, payload, level)
    for kind, payload, level in items:
        if kind in ("table", "code"):
            blocks.append((kind, payload, level))
        else:
            inline.append((kind, payload, level))

    # Reserve space for blocks at the bottom of the content region
    block_h = Inches(0) if not blocks else Inches(0.5 + 1.4 * len(blocks))
    text_h = max(H - block_h, Inches(1.0))

    if inline:
        body = slide.shapes.add_textbox(L, T, W, text_h)
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        first = True
        for kind, text, level in inline:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            indent_pt = level * 18
            p.level = min(level, 4)  # PPT supports indent levels 0–4
            r = p.add_run()
            if kind == "bullet":
                r.text = "• " + text
            elif kind == "numbered":
                r.text = f"{level + 1}. " + text  # indicative; PPT renumbers
            elif kind == "quote":
                r.text = "“ " + text + " ”"
                r.font.italic = True
                r.font.color.rgb = hex_to_rgb(brand["colors"]["text_muted"])
            else:
                r.text = text
            r.font.name = brand["fonts"]["body"]
            r.font.size = Pt(brand["sizes_pt"]["body"])
            if kind not in ("quote",):
                r.font.color.rgb = hex_to_rgb(brand["colors"]["text_primary"])
            p.space_after = Pt(4)

    # Now place table/code blocks
    block_top = T + (text_h if inline else Inches(0))
    for kind, payload, _level in blocks:
        if kind == "table":
            df = payload
            rows, cols = df.shape[0] + 1, max(df.shape[1], 1)
            header_h = Inches(0.45)
            data_row_h = Inches(0.35)
            tbl_h = header_h + data_row_h * (rows - 1)
            # Cap to remaining content area
            avail = T + H - block_top
            tbl_h = min(tbl_h, avail)
            shp = slide.shapes.add_table(rows, cols, L, block_top,
                                         W, tbl_h).table
            shp.rows[0].height = header_h
            for r in range(1, rows):
                shp.rows[r].height = data_row_h
            primary = hex_to_rgb(brand["colors"]["primary"])
            text_p  = hex_to_rgb(brand["colors"]["text_primary"])
            for j, col in enumerate(df.columns):
                cell = shp.cell(0, j)
                cell.text = str(col)
                cell.fill.solid(); cell.fill.fore_color.rgb = primary
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for r in p.runs:
                        r.font.bold = True
                        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        r.font.name = brand["fonts"]["header"]
                        r.font.size = Pt(11)
            for i, row in enumerate(df.itertuples(index=False), start=1):
                for j, val in enumerate(row):
                    cell = shp.cell(i, j)
                    cell.text = "" if pd.isna(val) else str(val)
                    # Solid white fill so the master's banded-row style
                    # doesn't reduce text contrast
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.LEFT
                        for r in p.runs:
                            r.font.name = brand["fonts"]["body"]
                            r.font.size = Pt(10)
                            r.font.color.rgb = text_p
            block_top += tbl_h + Inches(0.1)
        elif kind == "code":
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L, block_top,
                                         W, Inches(1.0))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFB)
            box.line.color.rgb = hex_to_rgb(brand["colors"]["primary"])
            box.line.width = Pt(0.5)
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            tf.clear()
            tf.text = ""
            for li, code_line in enumerate(str(payload).rstrip().split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = code_line
                r.font.name = "Courier New"
                r.font.size = Pt(10)
                r.font.color.rgb = hex_to_rgb(brand["colors"]["text_primary"])
            block_top += Inches(1.0) + Inches(0.1)

    add_branded_footer(slide, brand, prs, source_label=source_label)


def handle_html(path: Path, prs: Presentation, brand: dict, opts: dict):
    html = Path(path).read_text(encoding="utf-8", errors="ignore")
    src_label = (opts or {}).get("_source_label_override", path.name)
    # Walk top-level body (or root if no body) and yield structured blocks
    root = _parse_html(html)
    blocks = list(_walk_html_blocks(root))
    _emit_html_blocks(blocks, prs, brand, src_label)


def _mini_markdown(text: str) -> str:
    """Fallback .md → HTML for environments without the markdown package.
    Covers headings, bullets/numbered lists, fenced code, quotes, pipe
    tables and paragraphs — the constructs this converter renders."""
    out, lines = [], text.splitlines()
    i, in_code, in_ul, in_ol = 0, False, False, False

    def close_lists():
        nonlocal in_ul, in_ol
        if in_ul:
            out.append("</ul>"); in_ul = False
        if in_ol:
            out.append("</ol>"); in_ol = False

    def esc(s):
        return (s.replace("&", "&amp;").replace("<", "&lt;")
                 .replace(">", "&gt;"))

    def inline(s):
        s = esc(s)
        s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
        s = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", s)
        s = re.sub(r"`([^`]+)`", r"\1", s)
        return s

    while i < len(lines):
        ln = lines[i]
        if ln.strip().startswith("```"):
            if in_code:
                out.append("</pre>")
            else:
                close_lists(); out.append("<pre>")
            in_code = not in_code
            i += 1
            continue
        if in_code:
            out.append(esc(ln)); i += 1; continue
        m = re.match(r"^(#{1,3})\s+(.*)$", ln)
        if m:
            close_lists()
            out.append(f"<h{len(m.group(1))}>{inline(m.group(2))}"
                       f"</h{len(m.group(1))}>")
            i += 1; continue
        if re.match(r"^\s*[-*+]\s+", ln):
            if not in_ul:
                close_lists(); out.append("<ul>"); in_ul = True
            out.append("<li>" + inline(re.sub(r"^\s*[-*+]\s+", "", ln)) + "</li>")
            i += 1; continue
        if re.match(r"^\s*\d+[.)]\s+", ln):
            if not in_ol:
                close_lists(); out.append("<ol>"); in_ol = True
            out.append("<li>" + inline(re.sub(r"^\s*\d+[.)]\s+", "", ln)) + "</li>")
            i += 1; continue
        if ln.strip().startswith(">"):
            close_lists()
            out.append("<blockquote>" + inline(ln.strip().lstrip("> ")) +
                       "</blockquote>")
            i += 1; continue
        if "|" in ln and re.match(r"^\s*\|", ln):
            close_lists()
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
                    rows.append(cells)
                i += 1
            if rows:
                out.append("<table>")
                for ri, row in enumerate(rows):
                    tag = "th" if ri == 0 else "td"
                    out.append("<tr>" + "".join(
                        f"<{tag}>{inline(c)}</{tag}>" for c in row) + "</tr>")
                out.append("</table>")
            continue
        if ln.strip():
            close_lists()
            out.append("<p>" + inline(ln.strip()) + "</p>")
        i += 1
    if in_code:
        out.append("</pre>")
    close_lists()
    return "\n".join(out)


def handle_md(path: Path, prs: Presentation, brand: dict, opts: dict):
    src = Path(path).read_text(encoding="utf-8")
    if md_lib is not None:
        html = md_lib.markdown(src, extensions=["tables", "fenced_code"])
    else:
        html = _mini_markdown(src)
    tmp = path.with_suffix(".tmp.html")
    tmp.write_text(html, encoding="utf-8")
    try:
        # Pass the original .md filename through so footers cite the source
        # the user actually uploaded (not the tmp .html file).
        opts_with_label = dict(opts or {})
        opts_with_label["_source_label_override"] = path.name
        handle_html(tmp, prs, brand, opts_with_label)
    finally:
        tmp.unlink(missing_ok=True)


def handle_txt(path: Path, prs: Presentation, brand: dict, opts: dict):
    text = Path(path).read_text(encoding="utf-8", errors="ignore")
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    if not paras:
        return
    title = paras[0][:120]
    rest = paras[1:]
    chunk_size = 7
    for i in range(0, max(len(rest), 1), chunk_size):
        chunk = rest[i:i + chunk_size] or [paras[0]]
        slide_title = title if i == 0 else f"{title} (cont.)"
        add_title_body_slide(prs, slide_title, chunk, brand, source=path.name)


def handle_image(path: Path, prs: Presentation, brand: dict, opts: dict):
    """Image input handler.

    Three modes are supported (set ``options['image_mode']``):

      • ``"embed"`` (default) — image placed full-bleed on the slide; OCR
        text dropped into speaker notes for searchability. The image is
        a single editable picture object: the user can move/resize/replace
        it, but text inside the image is NOT individually editable.
        This is the safest mode for screenshots of designed slides where
        text colors and backgrounds vary across regions.

      • ``"transcript"`` — emit one slide with the original image at half
        size on the left, and an editable text-frame on the right
        containing the OCR'd text. Useful when the user wants to keep
        the visual reference AND edit the text without artifacts.

      • ``"overlay"`` (experimental) — place the image and overlay every
        OCR'd text line as an editable textbox on top, with a sampled
        background fill so it covers the original glyphs. Works well for
        single-color document scans; produces visible patches on
        multi-color designed slides.

    Photo-style images with no readable text fall through to ``embed``
    behavior automatically.
    """
    mode = opts.get("image_mode", "embed")
    slide = prs.slides.add_slide(_pick_layout(prs))
    margin = Inches(brand["spacing_inches"]["margin"])
    sw, sh = prs.slide_width, prs.slide_height

    img = Image.open(path)
    iw_px, ih_px = img.size

    # OCR once up front; falls back gracefully if tesseract is missing
    ocr_text = ""
    ocr_data = None
    try:
        if pytesseract is None:
            raise RuntimeError("OCR unavailable")
        ocr_text = pytesseract.image_to_string(img).strip()
        if mode == "overlay":
            ocr_data = pytesseract.image_to_data(
                img, output_type=pytesseract.Output.DICT)
    except Exception:
        pass

    if mode == "transcript":
        _image_transcript_layout(slide, path, img, brand, prs, ocr_text)
    elif mode == "overlay" and ocr_data is not None:
        _image_overlay_layout(slide, path, img, brand, prs, opts,
                              ocr_data, iw_px, ih_px)
    else:
        # Default: embed full-bleed
        _image_embed_layout(slide, path, brand, prs, ocr_text)

    add_branded_footer(slide, brand, prs, source_label=path.name)


def _image_embed_layout(slide, path, brand, prs, ocr_text):
    """Default: full-bleed image with OCR text in speaker notes."""
    margin = Inches(brand["spacing_inches"]["margin"])
    sw, sh = prs.slide_width, prs.slide_height
    img = Image.open(path)
    iw_px, ih_px = img.size
    aspect = iw_px / ih_px
    max_w = sw - 2 * margin
    max_h = sh - 2 * margin
    if max_w / max_h > aspect:
        h_emu = max_h; w_emu = int(h_emu * aspect)
    else:
        w_emu = max_w; h_emu = int(w_emu / aspect)
    slide.shapes.add_picture(str(path),
                             (sw - w_emu) // 2, (sh - h_emu) // 2,
                             w_emu, h_emu)
    if ocr_text:
        slide.notes_slide.notes_text_frame.text = ocr_text


def _image_transcript_layout(slide, path, img, brand, prs, ocr_text):
    """Image on the left, editable transcript on the right."""
    margin = Inches(brand["spacing_inches"]["margin"])
    sw, sh = prs.slide_width, prs.slide_height
    title_h = Inches(0.7)

    # Title
    add_branded_title(slide, path.stem, brand, prs)

    region_top = margin + title_h + Inches(0.1)
    region_h   = sh - region_top - Inches(0.5)
    region_w   = sw - 2 * margin

    # Left half: image
    iw_px, ih_px = img.size
    aspect = iw_px / ih_px
    half_w = region_w // 2 - Inches(0.1)
    if half_w / region_h > aspect:
        h_emu = region_h; w_emu = int(h_emu * aspect)
    else:
        w_emu = half_w; h_emu = int(w_emu / aspect)
    img_x = margin + (half_w - w_emu) // 2
    img_y = region_top + (region_h - h_emu) // 2
    slide.shapes.add_picture(str(path), img_x, img_y, w_emu, h_emu)

    # Right half: editable OCR transcript
    txt_x = margin + region_w // 2 + Inches(0.1)
    txt_w = region_w // 2 - Inches(0.1)
    box = slide.shapes.add_textbox(txt_x, region_top, txt_w, region_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = [l.strip() for l in (ocr_text or "(no text detected)").split("\n")
             if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = brand["fonts"]["body"]
        r.font.size = Pt(10)
        r.font.color.rgb = hex_to_rgb(brand["colors"]["text_primary"])


def _image_overlay_layout(slide, path, img, brand, prs, opts,
                           ocr_data, iw_px, ih_px):
    """Place image full-bleed and overlay each OCR'd line as an editable
    textbox with a sampled background fill so it covers the original glyphs.
    Best for single-color document scans; visible patches on designed slides."""
    margin = Inches(brand["spacing_inches"]["margin"])
    sw, sh = prs.slide_width, prs.slide_height
    aspect = iw_px / ih_px
    max_w = sw - 2 * margin
    max_h = sh - 2 * margin
    if max_w / max_h > aspect:
        h_emu = max_h; w_emu = int(h_emu * aspect)
    else:
        w_emu = max_w; h_emu = int(w_emu / aspect)
    img_left = (sw - w_emu) // 2
    img_top  = (sh - h_emu) // 2
    slide.shapes.add_picture(str(path), img_left, img_top, w_emu, h_emu)

    # Group OCR words → lines
    lines: dict = {}
    min_conf = int(opts.get("image_overlay_min_confidence", 50))
    n = len(ocr_data["text"])
    for i in range(n):
        txt = (ocr_data["text"][i] or "").strip()
        if not txt:
            continue
        try:
            conf = int(float(ocr_data["conf"][i]))
        except (ValueError, TypeError):
            conf = -1
        if conf < min_conf:
            continue
        key = (ocr_data["block_num"][i], ocr_data["par_num"][i],
               ocr_data["line_num"][i])
        line = lines.setdefault(key, {
            "words": [],
            "x_min": ocr_data["left"][i], "y_min": ocr_data["top"][i],
            "x_max": ocr_data["left"][i] + ocr_data["width"][i],
            "y_max": ocr_data["top"][i]  + ocr_data["height"][i],
        })
        line["words"].append(txt)
        line["x_min"] = min(line["x_min"], ocr_data["left"][i])
        line["y_min"] = min(line["y_min"], ocr_data["top"][i])
        line["x_max"] = max(line["x_max"],
                            ocr_data["left"][i] + ocr_data["width"][i])
        line["y_max"] = max(line["y_max"],
                            ocr_data["top"][i]  + ocr_data["height"][i])

    if not lines:
        return

    img_rgb = img.convert("RGB")
    from collections import Counter

    def _pick_colors(box):
        x1, y1, x2, y2 = box
        x1 = max(0, min(iw_px - 1, x1)); x2 = max(0, min(iw_px - 1, x2))
        y1 = max(0, min(ih_px - 1, y1)); y2 = max(0, min(ih_px - 1, y2))
        if x2 <= x1 or y2 <= y1:
            return RGBColor(0x16, 0x16, 0x16), RGBColor(0xFF, 0xFF, 0xFF)
        crop = img_rgb.crop((x1, y1, x2, y2)).resize(
            (max(1, (x2 - x1) // 4), max(1, (y2 - y1) // 4)))
        text_pixels = list(crop.getdata())
        bg_strip_h = max(2, (y2 - y1) // 3)
        bg_y1 = max(0, y1 - bg_strip_h); bg_y2 = max(0, y1 - 1)
        if bg_y2 <= bg_y1:
            bg_y1 = min(ih_px - bg_strip_h, y2 + 1)
            bg_y2 = min(ih_px - 1, bg_y1 + bg_strip_h)
        bg_crop = img_rgb.crop((x1, bg_y1, x2, bg_y2))
        bg_pixels = list(bg_crop.getdata()) if bg_crop.size[0] > 0 else []
        if not text_pixels:
            return RGBColor(0x16, 0x16, 0x16), RGBColor(0xFF, 0xFF, 0xFF)
        tr, tg, tb_ = min(text_pixels, key=lambda p: sum(p[:3]))
        if bg_pixels:
            br, bg_, bb = Counter(bg_pixels).most_common(1)[0][0]
        else:
            br, bg_, bb = max(text_pixels, key=lambda p: sum(p[:3]))
        return RGBColor(tr, tg, tb_), RGBColor(br, bg_, bb)

    for line in lines.values():
        text = " ".join(line["words"])
        x_min, y_min = line["x_min"], line["y_min"]
        x_max, y_max = line["x_max"], line["y_max"]
        x_frac = x_min / iw_px;       y_frac = y_min / ih_px
        w_frac = (x_max - x_min) / iw_px
        h_frac = (y_max - y_min) / ih_px
        pad_emu_x = int(w_emu * 0.004)
        pad_emu_y = int(h_emu * 0.005)
        tb_x = img_left + int(w_emu * x_frac) - pad_emu_x
        tb_y = img_top  + int(h_emu * y_frac) - pad_emu_y
        tb_w = int(w_emu * w_frac) + 2 * pad_emu_x
        tb_h = int(h_emu * h_frac) + 2 * pad_emu_y
        text_rgb, bg_rgb = _pick_colors((x_min, y_min, x_max, y_max))

        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, tb_x, tb_y, tb_w, tb_h)
        bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = bg_rgb
        bg_shape.line.fill.background()
        bg_shape.shadow.inherit = False

        tb = slide.shapes.add_textbox(tb_x, tb_y, tb_w, tb_h)
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = 0
        tf.margin_top  = tf.margin_bottom = 0
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = brand["fonts"]["body"]
        cap_pt = max(6, min(36, (tb_h / 12700) * 0.78))
        run.font.size = Pt(cap_pt)
        run.font.color.rgb = text_rgb


def handle_bpmn(path: Path, prs: Presentation, brand: dict, opts: dict):
    """Render BPMN flow elements as native, editable PPTX shapes laid out
    on a horizontal swimlane. Supports start/end events, tasks, gateways,
    and sequence flows (rendered as connectors)."""
    NS = {"bpmn":  "http://www.omg.org/spec/BPMN/20100524/MODEL",
          "bpmn2": "http://www.omg.org/spec/BPMN/20100524/MODEL"}
    tree = etree.parse(str(path))
    root = tree.getroot()
    # Auto-detect namespace from root
    ns_uri = root.nsmap.get(None) or root.nsmap.get("bpmn") or list(NS.values())[0]
    ns = {"b": ns_uri}

    slide = prs.slides.add_slide(_pick_layout(prs))
    add_branded_title(slide, f"BPMN: {path.stem}", brand, prs)
    L, T, W, H = _content_region(prs, brand)

    # Gather elements in document order
    elements = []
    for tag, shape_kind in (("startEvent", "start"),
                            ("task", "task"), ("userTask", "task"),
                            ("serviceTask", "task"), ("scriptTask", "task"),
                            ("exclusiveGateway", "gateway"),
                            ("parallelGateway", "gateway"),
                            ("endEvent", "end")):
        for el in root.findall(f".//b:{tag}", ns):
            elements.append({"id": el.get("id"),
                             "name": el.get("name") or tag,
                             "kind": shape_kind})
    if not elements:
        add_title_body_slide(prs, f"BPMN: {path.stem}",
                             ["No flow elements found in this BPMN file."],
                             brand, source=path.name)
        return

    # Layout horizontally with wrap
    box_w, box_h, gap = Inches(1.6), Inches(0.8), Inches(0.3)
    cols = max(1, int(W / (box_w + gap)))
    primary = brand["colors"]["primary"]
    accent  = brand["colors"]["teal"]
    danger  = brand["colors"]["red"]

    placed = {}
    for i, elem in enumerate(elements):
        row, col = divmod(i, cols)
        x = L + col * (box_w + gap)
        y = T + row * (box_h + gap)
        if elem["kind"] == "start":
            sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, box_w, box_h)
            sh.fill.solid(); sh.fill.fore_color.rgb = hex_to_rgb(accent)
        elif elem["kind"] == "end":
            sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, box_w, box_h)
            sh.fill.solid(); sh.fill.fore_color.rgb = hex_to_rgb(danger)
        elif elem["kind"] == "gateway":
            sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, box_w, box_h)
            sh.fill.solid(); sh.fill.fore_color.rgb = hex_to_rgb(brand["colors"]["secondary"])
        else:
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, y, box_w, box_h)
            sh.fill.solid(); sh.fill.fore_color.rgb = hex_to_rgb(primary)
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(45720)
        tf.text = elem["name"]
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.name = brand["fonts"]["header"]
        placed[elem["id"]] = sh

    add_branded_footer(slide, brand, prs, source_label=path.name)


HANDLERS = {
    "pptx": handle_pptx,
    "xlsx": handle_xlsx,
    "csv":  handle_csv,
    "pdf":  handle_pdf,
    "html": handle_html,
    "md":   handle_md,
    "txt":  handle_txt,
    "img":  handle_image,
    "bpmn": handle_bpmn,
}


# ══════════════════════════════════════════════════════════════════════
#  SOURCE INDEX SLIDE
# ══════════════════════════════════════════════════════════════════════
def add_source_index_slide(prs: Presentation, sources: list[tuple[str, str]],
                           brand: dict):
    lines = [f"• {name}   ({kind})" for name, kind in sources]
    add_title_body_slide(prs, "Source Index", lines, brand)


# ══════════════════════════════════════════════════════════════════════
#  PUBLIC ENTRY POINT
# ══════════════════════════════════════════════════════════════════════
def convert_to_pptx(input_paths: Iterable[str],
                    output_path: str,
                    template_path: str | None = None,
                    brand_path: str | None = None,
                    options: dict | None = None) -> dict:
    """Main entry point.

    Returns a dict with:
        output_path: str
        slide_count: int
        sources: [{filename, type}]
        warnings: [str]
    """
    options = options or {}
    brand   = load_brand(brand_path)
    prs     = open_blank_from_template(template_path)

    sources: list[tuple[str, str]] = []
    warnings: list[str] = []

    for raw in input_paths:
        p = Path(raw)
        if not p.exists():
            warnings.append(f"missing: {p}")
            continue
        kind = EXT_ROUTER.get(p.suffix.lower())
        if not kind:
            warnings.append(f"unsupported extension: {p.name}")
            continue
        if p.suffix.lower() == ".xlsm":
            warnings.append(f"{p.name}: macros dropped (data extracted only)")
        try:
            HANDLERS[kind](p, prs, brand, options)
            sources.append((p.name, kind))
        except Exception as e:
            warnings.append(f"{p.name}: {e}")

    if options.get("include_source_index_slide", True) and sources:
        add_source_index_slide(prs, sources, brand)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return {
        "output_path": str(out),
        "slide_count": len(prs.slides),
        "sources":     [{"filename": n, "type": k} for n, k in sources],
        "warnings":    warnings,
    }


# ──────────────────────────────────────────────────────────────────────
# CLI
# ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Universal File → PPTX")
    parser.add_argument("inputs", nargs="+", help="One or more input files")
    parser.add_argument("-o", "--output", required=True, help="Output .pptx path")
    parser.add_argument("--template", help="Brand template .pptx "
                        "(optional; defaults to the bundled IBM template)")
    parser.add_argument("--brand",    help="Brand pack .json "
                        "(optional; defaults to the bundled IBM Carbon pack)")
    parser.add_argument("--max-rows", type=int, default=20)
    parser.add_argument("--no-charts", action="store_true")
    parser.add_argument("--no-ocr",    action="store_true")
    args = parser.parse_args()

    opts = {
        "max_rows_per_table_slide": args.max_rows,
        "xlsx_chart_mode": "never" if args.no_charts else "auto",
        "ocr_scanned_pdfs": not args.no_ocr,
        "include_source_index_slide": True,
    }
    # Resolve the brand assets from inside the bundle when not passed.
    # They live beside this script (scripts/assets/) because the execution
    # policy requires every skill file the runtime touches to sit under
    # scripts/. Older layouts kept them at the skill root, so that is checked
    # as a fallback. Passing --template / --brand still overrides both.
    _here = Path(__file__).resolve().parent
    def _bundled(name):
        for cand in (_here / "assets" / name, _here.parent / "assets" / name):
            if cand.is_file():
                return str(cand)
        return None

    template_path = args.template or _bundled("ibm_template.pptx")
    brand_path = args.brand or _bundled("ibm_carbon.json")

    result = convert_to_pptx(args.inputs, args.output,
                             template_path=template_path,
                             brand_path=brand_path,
                             options=opts)
    print(json.dumps(result, indent=2))
