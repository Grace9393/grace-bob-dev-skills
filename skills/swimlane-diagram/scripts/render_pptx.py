#!/usr/bin/env python3
"""Render a swimlane diagram from canonical JSON to a PowerPoint slide (.pptx).

The output is fully editable: lanes are rectangles, nodes are standard
autoshapes, edges are connectors. Open in PowerPoint and tweak freely.

Usage:
    python render_pptx.py --input diagram.json --output diagram.pptx

Requires: python-pptx  (pip install python-pptx --break-system-packages)
"""

import argparse
import json
import sys
from collections import defaultdict, deque

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    sys.stderr.write(
        "python-pptx is required. Install with:\n"
        "    pip install python-pptx --break-system-packages\n"
    )
    sys.exit(1)


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

LANE_PALETTE = [
    ("E8F0FE", "4285F4"),
    ("FCE8E6", "EA4335"),
    ("FEF7E0", "FBBC04"),
    ("E6F4EA", "34A853"),
    ("F3E8FD", "A142F4"),
    ("E0F2F1", "00897B"),
    ("FFF3E0", "F57C00"),
    ("ECEFF1", "546E7A"),
]

NODE_W_IN = 1.5
NODE_H_IN = 0.7
COL_GAP_IN = 0.4
ROW_GAP_IN = 0.3
LANE_HEADER_W_IN = 1.2
LANE_HEADER_H_IN = 0.5
TITLE_H_IN = 0.5


def topo_rank(nodes, edges):
    node_ids = {n["id"] for n in nodes}
    outgoing = defaultdict(list)
    incoming = defaultdict(list)
    for e in edges:
        if e["from"] in node_ids and e["to"] in node_ids:
            outgoing[e["from"]].append(e["to"])
            incoming[e["to"]].append(e["from"])
    rank = {}
    indeg = {n["id"]: len(incoming[n["id"]]) for n in nodes}
    q = deque()
    for n in nodes:
        if indeg[n["id"]] == 0:
            rank[n["id"]] = 0
            q.append(n["id"])
    while len(rank) < len(nodes):
        while q:
            u = q.popleft()
            for v in outgoing[u]:
                if v in rank:
                    continue
                indeg[v] -= 1
                if indeg[v] == 0:
                    rank[v] = max((rank[p] for p in incoming[v] if p in rank), default=-1) + 1
                    q.append(v)
        if len(rank) < len(nodes):
            cands = [
                (indeg[n["id"]], i, n["id"])
                for i, n in enumerate(nodes)
                if n["id"] not in rank
            ]
            cands.sort()
            _, _, victim = cands[0]
            rank[victim] = max(
                (rank[p] for p in incoming[victim] if p in rank),
                default=max(rank.values(), default=-1),
            ) + 1
            q.append(victim)
    return rank


def shape_for_type(t):
    return {
        "start": MSO_SHAPE.ROUNDED_RECTANGLE,
        "end": MSO_SHAPE.ROUNDED_RECTANGLE,
        "process": MSO_SHAPE.RECTANGLE,
        "decision": MSO_SHAPE.DIAMOND,
        "subprocess": MSO_SHAPE.RECTANGLE,
        "data": MSO_SHAPE.PARALLELOGRAM,
        "document": MSO_SHAPE.FLOWCHART_DOCUMENT,
        "delay": MSO_SHAPE.FLOWCHART_DELAY,
        "annotation": MSO_SHAPE.RECTANGLE,
    }.get(t, MSO_SHAPE.RECTANGLE)


def fill_colors_for(t):
    """(fill_hex, line_hex) for a node type."""
    return {
        "start": ("E8F5E9", "43A047"),
        "end": ("FFEBEE", "E53935"),
        "decision": ("FFF8E1", "F9A825"),
        "data": ("FFFFFF", "9E9E9E"),
        "document": ("FFFFFF", "9E9E9E"),
        "delay": ("FFFFFF", "9E9E9E"),
        "annotation": ("FFFFFF", "9E9E9E"),
    }.get(t, ("FFFFFF", "546E7A"))


def set_fill(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)


def set_line(shape, hex_color, width_pt=1.0):
    shape.line.color.rgb = RGBColor.from_string(hex_color)
    shape.line.width = Pt(width_pt)


def set_text(shape, text, size_pt=11, bold=False, color_hex="212121", align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color_hex)


def render(data, out_path):
    title = data.get("title", "Swimlane Diagram")
    orientation = data.get("orientation", "horizontal")
    lanes = data["lanes"]
    nodes = data["nodes"]
    edges = data.get("edges", [])

    rank = topo_rank(nodes, edges)

    # De-dup same (lane, rank) cells by shifting
    cell_nodes = defaultdict(list)
    for n in nodes:
        cell_nodes[(n["lane"], rank[n["id"]])].append(n["id"])
    for (_, r), ids in cell_nodes.items():
        if len(ids) > 1:
            for off, nid in enumerate(ids[1:], start=1):
                rank[nid] = r + off
    max_rank = max(rank.values(), default=0)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), SLIDE_W - Inches(0.6), Inches(TITLE_H_IN)
    )
    set_text(title_box, title, size_pt=20, bold=True, color_hex="212121")

    lane_index = {lane["id"]: i for i, lane in enumerate(lanes)}

    # Sizing (work in plain inch floats, convert to Inches() at use sites)
    slide_w_in = SLIDE_W / 914400  # EMU to inches
    slide_h_in = SLIDE_H / 914400
    if orientation == "horizontal":
        usable_w = slide_w_in - LANE_HEADER_W_IN - 0.4
        cols = max_rank + 1
        col_w = max(NODE_W_IN + COL_GAP_IN, usable_w / max(cols, 1))
        usable_h = slide_h_in - TITLE_H_IN - 0.4
        row_h = max(NODE_H_IN + ROW_GAP_IN * 2, usable_h / max(len(lanes), 1))
    else:
        usable_w = slide_w_in - 0.4
        col_w = max(NODE_W_IN + COL_GAP_IN * 2, usable_w / max(len(lanes), 1))
        usable_h = slide_h_in - TITLE_H_IN - LANE_HEADER_H_IN - 0.4
        rows = max_rank + 1
        row_h = max(NODE_H_IN + ROW_GAP_IN, usable_h / max(rows, 1))

    positions = {}  # node_id -> (left_emu, top_emu, w_emu, h_emu)
    node_w_emu = Inches(NODE_W_IN)
    node_h_emu = Inches(NODE_H_IN)

    # Draw lanes
    if orientation == "horizontal":
        for i, lane in enumerate(lanes):
            fill, stroke = LANE_PALETTE[i % len(LANE_PALETTE)]
            top = Inches(TITLE_H_IN + 0.1) + Inches(i * row_h)
            # body
            body = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(LANE_HEADER_W_IN),
                top,
                SLIDE_W - Inches(LANE_HEADER_W_IN + 0.2),
                Inches(row_h),
            )
            set_fill(body, fill)
            set_line(body, stroke, 0.75)
            body.fill.transparency = 0  # solid, color light enough
            body.text_frame.text = ""
            # header
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.2),
                top,
                Inches(LANE_HEADER_W_IN - 0.2),
                Inches(row_h),
            )
            set_fill(header, fill)
            set_line(header, stroke, 1.0)
            set_text(header, lane["name"], size_pt=13, bold=True, color_hex=stroke)
    else:
        for i, lane in enumerate(lanes):
            fill, stroke = LANE_PALETTE[i % len(LANE_PALETTE)]
            left = Inches(0.2) + Inches(i * col_w)
            top_header = Inches(TITLE_H_IN + 0.1)
            top_body = top_header + Inches(LANE_HEADER_H_IN)
            body_h = SLIDE_H - top_body - Inches(0.2)
            body = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top_body, Inches(col_w), body_h
            )
            set_fill(body, fill)
            set_line(body, stroke, 0.75)
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top_header, Inches(col_w), Inches(LANE_HEADER_H_IN)
            )
            set_fill(header, fill)
            set_line(header, stroke, 1.0)
            set_text(header, lane["name"], size_pt=13, bold=True, color_hex=stroke)

    # Draw nodes
    node_shapes = {}
    for n in nodes:
        r = rank[n["id"]]
        li = lane_index[n["lane"]]
        if orientation == "horizontal":
            left = Inches(LANE_HEADER_W_IN + (col_w - NODE_W_IN) / 2 + r * col_w)
            top = Inches(TITLE_H_IN + 0.1 + li * row_h + (row_h - NODE_H_IN) / 2)
        else:
            left = Inches(0.2 + li * col_w + (col_w - NODE_W_IN) / 2)
            top = Inches(TITLE_H_IN + LANE_HEADER_H_IN + 0.1 + r * row_h + (row_h - NODE_H_IN) / 2)

        shape_type = shape_for_type(n.get("type", "process"))
        shape = slide.shapes.add_shape(shape_type, left, top, node_w_emu, node_h_emu)
        fill_hex, line_hex = fill_colors_for(n.get("type", "process"))
        set_fill(shape, fill_hex)
        set_line(shape, line_hex, 1.25)
        set_text(shape, n.get("label", ""), size_pt=11)
        positions[n["id"]] = (left, top, node_w_emu, node_h_emu)
        node_shapes[n["id"]] = shape

    # Draw edges as elbow connectors
    for e in edges:
        if e["from"] not in positions or e["to"] not in positions:
            continue
        a, b = node_shapes[e["from"]], node_shapes[e["to"]]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW,
            a.left + a.width,
            a.top + a.height // 2,
            b.left,
            b.top + b.height // 2,
        )
        connector.line.color.rgb = RGBColor.from_string("546E7A")
        connector.line.width = Pt(1.25)
        # Attach connector endpoints to shapes when possible
        try:
            connector.begin_connect(a, 3)  # right-center
            connector.end_connect(b, 1)    # left-center
        except Exception:
            pass
        if e.get("label"):
            lbl_box = slide.shapes.add_textbox(
                (a.left + a.width + b.left) // 2 - Inches(0.4),
                (a.top + a.height // 2 + b.top + b.height // 2) // 2 - Inches(0.15),
                Inches(0.8),
                Inches(0.3),
            )
            set_text(lbl_box, str(e["label"]), size_pt=9, color_hex="424242")

    prs.save(out_path)


def main():
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--input", "-i", required=True)
    p.add_argument("--output", "-o", required=True)
    args = p.parse_args()
    with open(args.input) as f:
        data = json.load(f)
    render(data, args.output)
    print(f"Wrote {args.output}", file=sys.stderr)


if __name__ == "__main__":
    main()
